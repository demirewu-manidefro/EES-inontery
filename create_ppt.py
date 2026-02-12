import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(md_file, output_file):
    prs = Presentation()
    
    # Define a clean layout
    # Slide Master indices usually: 0=Title, 1=Title+Content, etc.
    
    with open(md_file, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    current_slide = None
    text_frame = None
    
    # Regex patterns
    re_h1 = re.compile(r'^#\s+(.+)$')
    re_h2 = re.compile(r'^##\s+(.+)$')
    re_h3 = re.compile(r'^###\s+(.+)$')
    re_bullet = re.compile(r'^-\s+(.+)$')
    re_image = re.compile(r'!\[Screenshot:\s*(.+)\]')

    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Title Slide (H1)
        match_h1 = re_h1.match(line)
        if match_h1:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = match_h1.group(1)
            subtitle.text = "Project Documentation"
            continue
            
        # New Slide (H2)
        match_h2 = re_h2.match(line)
        if match_h2:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = match_h2.group(1)
            
            # Reset text frame for new slide
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear() 
            continue
            
        # Subheader (H3) - Make bold paragraph
        match_h3 = re_h3.match(line)
        if match_h3 and text_frame:
            p = text_frame.add_paragraph()
            p.text = match_h3.group(1)
            p.font.bold = True
            p.font.size = Pt(20)
            p.space_before = Pt(12)
            continue
            
        # Bullet Points
        match_bullet = re_bullet.match(line)
        if match_bullet and text_frame:
            p = text_frame.add_paragraph()
            p.text = match_bullet.group(1)
            p.level = 0
            continue
            
        # Check for nested bullets (indent 4 spaces)
        if line.startswith('    - ') and text_frame:
            p = text_frame.add_paragraph()
            p.text = line.strip()[2:]
            p.level = 1
            continue

        # Image Placeholders
        match_img = re_image.search(line)
        if match_img and text_frame:
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f"[INSERT IMAGE: {match_img.group(1)}]"
            run.font.color.rgb = RGBColor(255, 0, 0)
            run.font.bold = True
            continue
            
        # Normal Text (if slide exists)
        if current_slide is None and not line.startswith('#'):
            # Text before first slide? Ignore or add to title
            pass
        elif text_frame and not line.startswith('!['):
            # Add as plain paragraph if it's not looking like other markdown syntax
            # Avoid adding every random line, but add significant text
            if len(line) > 3 and not line.startswith('---') and not line.startswith('>'):
                p = text_frame.add_paragraph()
                p.text = line
                p.level = 0

    prs.save(output_file)
    print(f"Successfully generated {output_file}")

if __name__ == "__main__":
    try:
        create_presentation('project_documentation.md', 'Employee_Management_System_Docs.pptx')
    except Exception as e:
        print(f"Error: {e}")
